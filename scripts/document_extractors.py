#!/usr/bin/env python3
"""Extract text and metadata from office documents for knowledge base ingest."""

from __future__ import annotations

import hashlib
import mimetypes
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".txt",
    ".md",
    ".csv",
}

SOURCE_TYPE_BY_EXTENSION = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".txt": "document",
    ".md": "document",
    ".csv": "spreadsheet",
}

MIME_BY_EXTENSION = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
}


def sha256_path(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def guess_mime_type(path: Path) -> str:
    return MIME_BY_EXTENSION.get(path.suffix.lower()) or mimetypes.guess_type(path.name)[0] or "application/octet-stream"


def guess_source_type(path: Path) -> str:
    return SOURCE_TYPE_BY_EXTENSION.get(path.suffix.lower(), "document")


def supported_path(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXTENSIONS


def extract_document_payload(path: Path) -> dict[str, object]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".xlsx":
        return _extract_xlsx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix in {".txt", ".md", ".csv"}:
        return _extract_text_file(path)
    raise ValueError(f"Unsupported document type: {path.suffix}")


def _base_payload(path: Path, extracted_text: str, normalized_text: str, metadata: dict) -> dict[str, object]:
    return {
        "external_id": f"local:{path.resolve()}",
        "checksum": sha256_path(path),
        "title": path.stem,
        "mime_type": guess_mime_type(path),
        "source_type": guess_source_type(path),
        "raw_storage_url": str(path.resolve()),
        "extracted_text": extracted_text[:60000],
        "normalized_text": normalized_text[:60000],
        "source_metadata": {
            "ingest_origin": "local_filesystem",
            "absolute_path": str(path.resolve()),
            "file_name": path.name,
            "extension": path.suffix.lower(),
            "file_size_bytes": path.stat().st_size,
        },
        "document_metadata": metadata,
    }


def _extract_pdf(path: Path) -> dict[str, object]:
    reader = PdfReader(str(path))
    page_texts = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            page_texts.append(f"Page {index}\n{text}")

    normalized_text = "\n\n".join(page_texts).strip()
    extracted_text = "\n".join(
        [
            f"File name: {path.name}",
            f"Page count: {len(reader.pages)}",
            "",
            normalized_text,
        ]
    ).strip()

    return _base_payload(
        path,
        extracted_text=extracted_text,
        normalized_text=normalized_text,
        metadata={
            "artifact_type": "pdf",
            "page_count": len(reader.pages),
        },
    )


def _extract_docx(path: Path) -> dict[str, object]:
    document = DocxDocument(str(path))
    sections = []
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    if paragraphs:
        sections.append("\n".join(paragraphs))

    table_rows = []
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_rows.append(" | ".join(cells))
    if table_rows:
        sections.append("Tables\n" + "\n".join(table_rows))

    normalized_text = "\n\n".join(section for section in sections if section).strip()
    extracted_text = "\n".join(
        [
            f"File name: {path.name}",
            "",
            normalized_text,
        ]
    ).strip()

    return _base_payload(
        path,
        extracted_text=extracted_text,
        normalized_text=normalized_text,
        metadata={
            "artifact_type": "docx",
            "paragraph_count": len(paragraphs),
            "table_count": len(document.tables),
        },
    )


def _extract_xlsx(path: Path) -> dict[str, object]:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    sheet_sections = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        if rows:
            sheet_sections.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))

    normalized_text = "\n\n".join(sheet_sections).strip()
    extracted_text = "\n".join(
        [
            f"File name: {path.name}",
            f"Sheet count: {len(workbook.sheetnames)}",
            "",
            normalized_text,
        ]
    ).strip()

    return _base_payload(
        path,
        extracted_text=extracted_text,
        normalized_text=normalized_text,
        metadata={
            "artifact_type": "xlsx",
            "sheet_names": workbook.sheetnames,
            "sheet_count": len(workbook.sheetnames),
        },
    )


def _extract_pptx(path: Path) -> dict[str, object]:
    presentation = Presentation(str(path))
    slide_sections = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                texts.append(text.strip())
        if texts:
            slide_sections.append(f"Slide {slide_index}\n" + "\n".join(texts))

    normalized_text = "\n\n".join(slide_sections).strip()
    extracted_text = "\n".join(
        [
            f"File name: {path.name}",
            f"Slide count: {len(presentation.slides)}",
            "",
            normalized_text,
        ]
    ).strip()

    return _base_payload(
        path,
        extracted_text=extracted_text,
        normalized_text=normalized_text,
        metadata={
            "artifact_type": "pptx",
            "slide_count": len(presentation.slides),
        },
    )


def _extract_text_file(path: Path) -> dict[str, object]:
    normalized_text = path.read_text(encoding="utf-8", errors="ignore").strip()
    extracted_text = "\n".join(
        [
            f"File name: {path.name}",
            "",
            normalized_text,
        ]
    ).strip()

    return _base_payload(
        path,
        extracted_text=extracted_text,
        normalized_text=normalized_text,
        metadata={
            "artifact_type": path.suffix.lower().lstrip(".") or "text",
        },
    )
